from fastapi import FastAPI, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse, HTMLResponse
from fastapi.staticfiles import StaticFiles
from datetime import datetime
import uvicorn
import io
import os
import requests
from pptx import Presentation
import shutil

app = FastAPI(title="AI PPT Generator API", version="2.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# 创建必要的目录
os.makedirs("templates", exist_ok=True)
os.makedirs("generated", exist_ok=True)

# ---------------- SiliconFlow LLM ----------------
SILICONFLOW_API_URL = "https://api.siliconflow.cn/v1/chat/completions"
SILICONFLOW_API_KEY = os.getenv("SILICONFLOW_API_KEY")

if not SILICONFLOW_API_KEY:
    SILICONFLOW_API_KEY = "your own apikey".strip()

def call_llm(prompt: str) -> str:
    if not SILICONFLOW_API_KEY or not SILICONFLOW_API_KEY.startswith("sk-"):
        raise ValueError("无效的 API KEY，请检查配置")
    headers = {
        "Authorization": f"Bearer {SILICONFLOW_API_KEY}",
        "Content-Type": "application/json",
    }
    body = {
        "model": "deepseek-ai/DeepSeek-V3",
        "messages": [{"role": "user", "content": prompt}],
        "temperature": 0.3,
    }
    resp = requests.post(SILICONFLOW_API_URL, headers=headers, json=body, timeout=60)
    if resp.status_code != 200:
        raise RuntimeError(f"调用 LLM 失败: {resp.status_code}, {resp.text}")
    data = resp.json()
    return data["choices"][0]["message"]["content"].strip()

# ---------------- PPT模板 ----------------
TEMPLATES = {
    "default": {
        "name": "默认模板",
        "description": "简洁专业的商务风格",
        "file": "templates/default.pptx",
        "cover_layout": 0,  # 封面页布局索引
        "content_layouts": [1, 2, 3]  # 三种内容页布局索引，按顺序循环
    },
    "blue": {
        "name": "蓝色商务",
        "description": "专业商务风格，蓝色主题",
        "file": "templates/blue.pptx",
        "cover_layout": 0,
        "content_layouts": [1, 2, 3]
    },
    "green": {
        "name": "绿色环保",
        "description": "清新环保风格，绿色主题",
        "file": "templates/green.pptx",
        "cover_layout": 0,
        "content_layouts": [1, 2, 3]
    },
    "red": {
        "name": "红色活力",
        "description": "活力四射风格，红色主题",
        "file": "templates/red.pptx",
        "cover_layout": 0,
        "content_layouts": [1, 2, 3]
    },
    "dark": {
        "name": "深色专业",
        "description": "深色背景专业风格",
        "file": "templates/dark.pptx",
        "cover_layout": 0,
        "content_layouts": [1, 2, 3]
    }
}

# ---------------- PPT生成 ----------------
def text_to_slides(title: str, outline_text: str, template_name: str):
    """使用外部模板文件生成PPT，内容页按顺序循环使用多种布局"""
    template = TEMPLATES.get(template_name, TEMPLATES["default"])
    template_file = template["file"]
    cover_layout_idx = template["cover_layout"]
    content_layouts = template["content_layouts"]
    
    # 检查模板文件是否存在
    if not os.path.exists(template_file):
        raise FileNotFoundError(f"模板文件不存在: {template_file}")
    
    # 从模板文件创建演示文稿
    prs = Presentation(template_file)
    
    # 验证布局索引有效性
    if cover_layout_idx >= len(prs.slide_layouts):
        raise ValueError(f"封面布局索引 {cover_layout_idx} 超出模板实际布局数量")
    for idx in content_layouts:
        if idx >= len(prs.slide_layouts):
            raise ValueError(f"内容布局索引 {idx} 超出模板实际布局数量")
    
    # 封面页 - 使用指定的封面布局
    slide_layout = prs.slide_layouts[cover_layout_idx]
    slide = prs.slides.add_slide(slide_layout)
    if slide.shapes.title:
        slide.shapes.title.text = title
    # 清空副标题（如果有）
    if len(slide.placeholders) > 1 and slide.placeholders[1]:
        slide.placeholders[1].text = ""
    
    # 内容页 - 按顺序循环使用不同布局
    lines = [l.strip() for l in outline_text.split("\n") if l.strip()]
    section_title = None
    bullets = []
    layout_counter = 0  # 布局循环计数器
    
    for line in lines + [""]:  # 加空行确保最后一页生成
        if not line:  # 处理空行触发页面生成
            if section_title and bullets:
                # 获取当前布局索引（循环）
                current_layout_idx = content_layouts[layout_counter % len(content_layouts)]
                slide_layout = prs.slide_layouts[current_layout_idx]
                slide = prs.slides.add_slide(slide_layout)
                
                # 设置标题和内容
                slide.shapes.title.text = section_title
                if len(slide.placeholders) > 1:
                    tf = slide.placeholders[1].text_frame
                    tf.clear()
                    for b in bullets:
                        p = tf.add_paragraph()
                        p.text = b
                
                layout_counter += 1  # 计数器自增
                section_title = None
                bullets = []
        elif not line.startswith("-"):  # 新标题
            if section_title and bullets:
                # 获取当前布局索引（循环）
                current_layout_idx = content_layouts[layout_counter % len(content_layouts)]
                slide_layout = prs.slide_layouts[current_layout_idx]
                slide = prs.slides.add_slide(slide_layout)
                
                # 设置标题和内容
                slide.shapes.title.text = section_title
                if len(slide.placeholders) > 1:
                    tf = slide.placeholders[1].text_frame
                    tf.clear()
                    for b in bullets:
                        p = tf.add_paragraph()
                        p.text = b
                
                layout_counter += 1  # 计数器自增
            section_title = line
            bullets = []
        else:
            bullets.append(line[1:].strip())
    
    return prs

# ---------------- API 接口 ----------------
@app.get("/api/templates")
async def get_templates():
    """获取可用模板列表"""
    return [{"id": k, **v} for k, v in TEMPLATES.items()]

@app.post("/api/outline")
async def generate_outline(keyword: str = Form(...)):
    """根据关键词生成 AI 标题和大纲"""
    try:
        # 生成标题
        prompt_title = f"请为以下主题生成一个简洁、吸引人的中文PPT标题：\n主题：{keyword}\n只输出标题文本，不要多余说明。"
        title = call_llm(prompt_title)

        # 生成大纲
        prompt_outline = f"""
请根据关键词「{keyword}」生成一个演示文稿大纲，要求：
• 3~5 个部分，每部分有一个标题和 2~4 个要点

• 用中文，每个部分之间空一行，每个要点用 '-' 开头

• 输出示例：

部分标题1
• 要点1

• 要点2


部分标题2
• 要点1

• 要点2

请严格按照此格式输出，便于程序解析。
"""
        outline = call_llm(prompt_outline)
        return {"title": title, "outline": outline}

    except Exception as e:
        return JSONResponse(status_code=500, content={"error": str(e)})

@app.post("/api/generate")
async def generate_ppt(
    title: str = Form(...), 
    content: str = Form(...),
    template: str = Form("default")
):
    """根据用户修改后的标题和内容生成 PPT"""
    if template not in TEMPLATES:
        return JSONResponse(status_code=400, content={"error": "无效的模板选择"})

    try:
        prs = text_to_slides(title, content, template)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        ts = datetime.now().strftime("%Y%m%d-%H%M%S")
        fname = f"PPT_{ts}.pptx"
        out_path = os.path.join("generated", fname)
        os.makedirs("generated", exist_ok=True)
        with open(out_path, "wb") as f:
            f.write(buf.read())

        return {"filename": fname, "url": f"/download/{fname}"}
    except Exception as e:
        return JSONResponse(status_code=500, content={"error": f"生成PPT失败: {str(e)}"})

@app.get("/download/{fname}")
async def download_file(fname: str):
    path = os.path.join("generated", fname)
    if not os.path.exists(path):
        return JSONResponse(status_code=404, content={"error": "文件不存在"})
    return FileResponse(
        path,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=fname,
    )

# 提供前端HTML文件
@app.get("/")
async def get_frontend():
    with open("index.html", "r", encoding="utf-8") as f:
        html_content = f.read()
    return HTMLResponse(content=html_content)

if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=8000)